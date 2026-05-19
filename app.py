"""
Big Bets PPT 解析 API + MCP 协议支持
MCP 工具接收智能体从知识库读取的 PPT 文本内容，做精确字段提取
"""

from flask import Flask, request, jsonify, Response
from pptx import Presentation
from io import BytesIO
import re
import json

app = Flask(__name__)

# 需要提取的字段关键词
FIELD_KEYWORDS = [
    "Key issues",
    "Insights to why",
    "Objective",
    "Project Boundary",
    "Project Scope",
    "Any Hypotheses",
    "Expected output",
]

# 8 个 Big Bets 名称
BIG_BETS_NAMES = [
    "Device Management",
    "Future Restaurant",
    "Store Agent",
    "MACE",
    "McCruiser",
    "Next Gen DMB",
    "Smart Production",
    "Digital Drive Through",
]


def parse_text_content(text_content):
    """
    从文本内容中解析 Big Bets 信息。
    智能体从知识库读取 PPT 后，把文本发过来，这里做精确的字段提取。
    """
    results = []

    # 按 Big Bet 名称分割文本
    # 先找到每个 Big Bet 在文本中的位置
    positions = []
    for name in BIG_BETS_NAMES:
        # 查找包含 Big Bet 名称和 BBO 的行
        pattern = re.compile(
            rf'({re.escape(name)})\s*[-–—]\s*BBO\s+(\w+)',
            re.IGNORECASE
        )
        for match in pattern.finditer(text_content):
            positions.append({
                "name": name,
                "owner": match.group(2),
                "start": match.start()
            })

    # 如果没找到标准格式，尝试更宽松的匹配
    if not positions:
        for name in BIG_BETS_NAMES:
            pattern = re.compile(
                rf'({re.escape(name)})',
                re.IGNORECASE
            )
            for match in pattern.finditer(text_content):
                # 检查附近是否有 BBO
                context = text_content[match.start():match.start() + 200]
                bbo_match = re.search(r'BBO\s+(\w+)', context, re.IGNORECASE)
                owner = bbo_match.group(1) if bbo_match else ""
                positions.append({
                    "name": name,
                    "owner": owner,
                    "start": match.start()
                })
                break  # 只取第一个匹配

    # 按位置排序
    positions.sort(key=lambda x: x["start"])

    # 为每个 Big Bet 提取对应的文本段
    for i, pos in enumerate(positions):
        # 确定这个 Big Bet 的文本范围
        start = pos["start"]
        end = positions[i + 1]["start"] if i + 1 < len(positions) else len(text_content)
        section_text = text_content[start:end]

        # 将文本按行分割
        lines = [line.strip() for line in section_text.split('\n') if line.strip()]

        # 提取各字段
        record = {
            "Big Bets": pos["name"],
            "Big Bets Owner": pos["owner"],
            "Key issues": extract_field_from_lines(lines, "Key issues"),
            "Insights to why": extract_field_from_lines(lines, "Insights to why"),
            "Objective": extract_field_from_lines(lines, "Objective"),
            "Project Boundary": extract_field_from_lines(lines, "Project Boundary"),
            "Project Scope": extract_field_from_lines(lines, "Project Scope"),
            "Any Hypotheses": extract_field_from_lines(lines, "Any Hypotheses"),
            "Expected output/Success": extract_field_from_lines(lines, "Expected output"),
        }

        results.append(record)

    return results


def extract_field_from_lines(lines, field_keyword):
    """从行列表中提取某个字段的内容"""
    start_index = -1
    keyword_lower = field_keyword.lower().replace(":", "").strip()

    for i, line in enumerate(lines):
        cleaned = line.lower().replace(":", "").replace("：", "").strip()
        if keyword_lower in cleaned:
            start_index = i
            break

    if start_index == -1:
        return ""

    content_parts = []

    # 检查关键词行本身是否包含内容
    keyword_line = lines[start_index]
    after_keyword = re.split(r'[:：]', keyword_line, maxsplit=1)
    if len(after_keyword) > 1 and after_keyword[1].strip():
        content_parts.append(after_keyword[1].strip())

    # 收集后续行直到下一个字段
    for i in range(start_index + 1, len(lines)):
        line = lines[i]
        is_next_field = False
        for kw in FIELD_KEYWORDS:
            cleaned = line.lower().replace(":", "").replace("：", "").strip()
            kw_cleaned = kw.lower().replace(":", "").strip()
            if kw_cleaned in cleaned and len(line) < len(kw) + 50:
                is_next_field = True
                break
        if is_next_field:
            break
        content_parts.append(line)

    return "\n".join(content_parts)


def parse_ppt(file_bytes):
    """解析 PPT 二进制文件"""
    prs = Presentation(BytesIO(file_bytes))
    results = []

    for slide in prs.slides:
        all_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        all_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            all_texts.append(text)

        if not all_texts:
            continue

        # 提取标题
        big_bet = ""
        owner = ""
        for text in all_texts:
            for name in BIG_BETS_NAMES:
                if name.lower() in text.lower():
                    big_bet = name
                    bbo_match = re.search(r'BBO\s+(.+)', text, re.IGNORECASE)
                    if bbo_match:
                        owner = bbo_match.group(1).strip()
                    break
            if big_bet:
                break

        if not big_bet:
            continue

        def find_field(field_keyword):
            start_index = -1
            for i, text in enumerate(all_texts):
                cleaned = text.lower().replace(":", "").replace("：", "").strip()
                kw_cleaned = field_keyword.lower().replace(":", "").strip()
                if kw_cleaned in cleaned:
                    start_index = i
                    break
            if start_index == -1:
                return ""
            content_parts = []
            keyword_line = all_texts[start_index]
            after_keyword = re.split(r'[:：]', keyword_line, maxsplit=1)
            if len(after_keyword) > 1 and after_keyword[1].strip():
                content_parts.append(after_keyword[1].strip())
            for i in range(start_index + 1, len(all_texts)):
                text = all_texts[i]
                is_next_field = False
                for kw in FIELD_KEYWORDS:
                    cleaned = text.lower().replace(":", "").replace("：", "").strip()
                    kw_cleaned = kw.lower().replace(":", "").strip()
                    if kw_cleaned in cleaned and len(text) < len(kw) + 30:
                        is_next_field = True
                        break
                if is_next_field:
                    break
                content_parts.append(text)
            return "\n".join(content_parts)

        record = {
            "Big Bets": big_bet,
            "Big Bets Owner": owner,
            "Key issues": find_field("Key issues"),
            "Insights to why": find_field("Insights to why"),
            "Objective": find_field("Objective"),
            "Project Boundary": find_field("Project Boundary"),
            "Project Scope": find_field("Project Scope"),
            "Any Hypotheses": find_field("Any Hypotheses"),
            "Expected output/Success": find_field("Expected output"),
        }
        results.append(record)

    return results


# ============ REST API ============

@app.route("/parse", methods=["POST"])
def parse_endpoint():
    if "file" in request.files:
        file_bytes = request.files["file"].read()
    else:
        file_bytes = request.get_data()
    if not file_bytes:
        return jsonify({"error": "No file provided"}), 400
    try:
        results = parse_ppt(file_bytes)
        return jsonify({"success": True, "count": len(results), "data": results})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


# ============ MCP 协议 ============

def make_sse_message(event_type, data):
    return f"event: {event_type}\ndata: {json.dumps(data)}\n\n"


@app.route("/mcp", methods=["GET", "POST"])
def mcp_endpoint():
    if request.method == "GET":
        def generate():
            yield make_sse_message("endpoint", "/mcp/messages")
        return Response(generate(), mimetype="text/event-stream",
                       headers={"Cache-Control": "no-cache", "Connection": "keep-alive"})

    msg = request.get_json()
    method = msg.get("method", "")
    msg_id = msg.get("id")

    if method == "initialize":
        return jsonify({
            "jsonrpc": "2.0",
            "id": msg_id,
            "result": {
                "protocolVersion": "2024-11-05",
                "capabilities": {"tools": {"listChanged": False}},
                "serverInfo": {"name": "BigBets Parser", "version": "2.0.0"}
            }
        })

    elif method == "tools/list":
        return jsonify({
            "jsonrpc": "2.0",
            "id": msg_id,
            "result": {
                "tools": [
                    {
                        "name": "parse_bigbets_text",
                        "description": "接收从PPT中读取的文本内容，精确提取每个Big Bet的9个字段（Big Bets、Big Bets Owner、Key issues、Insights to why、Objective、Project Boundary、Project Scope、Any Hypotheses、Expected output/Success）。请将PPT的完整文本内容作为text参数传入。",
                        "inputSchema": {
                            "type": "object",
                            "properties": {
                                "text": {
                                    "type": "string",
                                    "description": "从PPT文件中读取的完整文本内容。包含所有页面的文字，每个Big Bet的标题格式为'Big Bet名称 – BBO 人名'。"
                                }
                            },
                            "required": ["text"]
                        }
                    }
                ]
            }
        })

    elif method == "tools/call":
        tool_name = msg.get("params", {}).get("name", "")
        arguments = msg.get("params", {}).get("arguments", {})

        if tool_name == "parse_bigbets_text":
            text_content = arguments.get("text", "")
            if not text_content:
                return jsonify({
                    "jsonrpc": "2.0",
                    "id": msg_id,
                    "result": {
                        "content": [{"type": "text", "text": "错误：未提供文本内容。请从知识库中读取PPT的完整文本并传入text参数。"}]
                    }
                })

            try:
                results = parse_text_content(text_content)
                if not results:
                    return jsonify({
                        "jsonrpc": "2.0",
                        "id": msg_id,
                        "result": {
                            "content": [{"type": "text", "text": "未能从文本中提取到Big Bets信息。请确保传入的是完整的PPT文本内容，包含标题格式'Big Bet名称 – BBO 人名'。"}]
                        }
                    })

                result_text = json.dumps(results, ensure_ascii=False, indent=2)
                return jsonify({
                    "jsonrpc": "2.0",
                    "id": msg_id,
                    "result": {
                        "content": [{
                            "type": "text",
                            "text": f"成功提取 {len(results)} 个Big Bets的数据：\n\n{result_text}"
                        }]
                    }
                })
            except Exception as e:
                return jsonify({
                    "jsonrpc": "2.0",
                    "id": msg_id,
                    "result": {
                        "content": [{"type": "text", "text": f"解析失败: {str(e)}"}],
                        "isError": True
                    }
                })

        else:
            return jsonify({
                "jsonrpc": "2.0",
                "id": msg_id,
                "error": {"code": -32601, "message": f"Unknown tool: {tool_name}"}
            })

    elif method == "notifications/initialized":
        return "", 204

    else:
        return jsonify({
            "jsonrpc": "2.0",
            "id": msg_id,
            "error": {"code": -32601, "message": f"Method not found: {method}"}
        })


@app.route("/mcp/messages", methods=["POST"])
def mcp_messages():
    return mcp_endpoint()


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
